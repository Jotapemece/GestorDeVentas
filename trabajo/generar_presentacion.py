# -*- coding: utf-8 -*-
"""Genera Presentacion_Defensa_JuanMoya.pptx (17 diapositivas, defensa de pasantía).

Estilo: verde corporativo #2E5E3A + crema #E8E2D0, limpio y profesional.
Personalidad: kicker numerado por bloque, palabras clave resaltadas en verde,
callouts de frases de impacto y barra de progreso en el pie.
"""
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ----------------------------- Paleta -----------------------------
GREEN = RGBColor(0x2E, 0x5E, 0x3A)
CREAM = RGBColor(0xE8, 0xE2, 0xD0)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = 'Calibri'

RUTA_SALIDA = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'Presentacion_Defensa_JuanMoya.pptx')
RUTA_ISHIKAWA = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'imagenes', 'ishikawa_geo.png')

W = 13.333
H = 7.5
MARGIN = 0.5
CONTENT_W = W - 2 * MARGIN  # 12.333
TOTAL_SLIDES = 17


# ----------------------------- Helpers -----------------------------
def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_rich_runs(p, text, size, base_color=DARK, kw_color=GREEN, base_bold=False):
    """Añade el texto a `p` partiendo por `**kw**` (verde negrita)."""
    parts = re.split(r'\*\*(.+?)\*\*', text)
    for i, part in enumerate(parts):
        if not part:
            continue
        r = p.add_run()
        r.text = part
        r.font.size = Pt(size)
        r.font.name = FONT
        if i % 2 == 1:
            r.font.bold = True
            r.font.color.rgb = kw_color
        else:
            r.font.bold = base_bold
            r.font.color.rgb = base_color


def add_title(slide, text, kicker=None):
    """Título en esquina superior izquierda + línea de subrayado separada.

    Si `kicker` se indica, se añade una etiqueta numerada arriba del título.
    """
    if kicker:
        kt = slide.shapes.add_textbox(Inches(MARGIN), Inches(0.15), Inches(11.5), Inches(0.25))
        kp = kt.text_frame.paragraphs[0]
        kr = kp.add_run()
        kr.text = kicker.upper()
        kr.font.size = Pt(11)
        kr.font.bold = True
        kr.font.color.rgb = GREEN
        kr.font.name = FONT
        top_title = 0.45
        top_line = 1.12
    else:
        top_title = 0.32
        top_line = 1.0
    tb = slide.shapes.add_textbox(Inches(MARGIN), Inches(top_title), Inches(11.5), Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = GREEN
    r.font.name = FONT
    # Línea de subrayado, separada del texto (gap)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(top_line), Inches(CONTENT_W), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN
    line.line.fill.background()
    return tb


def add_footer(slide, n):
    # Barra de progreso fina
    shape(slide, MSO_SHAPE.RECTANGLE, MARGIN, 6.95, CONTENT_W, 0.05, fill=CREAM, line=None)
    fillw = CONTENT_W * (n / TOTAL_SLIDES)
    if fillw > 0.05:
        shape(slide, MSO_SHAPE.RECTANGLE, MARGIN, 6.95, fillw, 0.05, fill=GREEN, line=None)
    # Texto
    tb = slide.shapes.add_textbox(Inches(MARGIN), Inches(7.08), Inches(9), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = 'Juan Moya · Defensa de Pasantía'
    r.font.size = Pt(9)
    r.font.color.rgb = GRAY
    r.font.name = FONT
    nb = slide.shapes.add_textbox(Inches(W - 1.2), Inches(7.08), Inches(0.8), Inches(0.3))
    p = nb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f'{n:02d}'
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = GREEN
    r.font.name = FONT


def add_bullets(slide, items, top=1.35, left=MARGIN, width=CONTENT_W, height=5.4):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        text, level = it if isinstance(it, tuple) else (it, 0)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(9)
        p.level = level
        marker = '        ▪  ' if level > 0 else '▪  '
        rb = p.add_run()
        rb.text = marker
        rb.font.color.rgb = GREEN
        rb.font.bold = True
        rb.font.size = Pt(16 if level == 0 else 14)
        rb.font.name = FONT
        _add_rich_runs(p, text, 17 if level == 0 else 15, base_color=DARK, kw_color=GREEN, base_bold=False)


def shape(slide, kind, left, top, width, height, fill=None, line=GREEN, line_w=1, dash=False):
    sh = slide.shapes.add_shape(kind, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
        if dash:
            sh.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    sh.shadow.inherit = False
    return sh


def shape_text(sh, text, size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER):
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    kw = WHITE if color == GREEN else GREEN
    _add_rich_runs(p, text, size, base_color=color, kw_color=kw, base_bold=bold)


def pill(slide, left, top, width, height, text, fill=CREAM):
    sh = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill=fill, line=GREEN, line_w=1)
    shape_text(sh, text, size=12, bold=True, color=GREEN)
    return sh


def card(slide, left, top, width, height, title, body):
    sh = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill=CREAM, line=GREEN, line_w=1)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(6)
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = GREEN
    r.font.name = FONT
    p2 = tf.add_paragraph()
    _add_rich_runs(p2, body, 12, base_color=DARK, kw_color=GREEN, base_bold=False)
    return sh


def placeholder(slide, left, top, width, height, caption):
    sh = shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, fill=CREAM, line=GREEN, line_w=1.25, dash=True)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = caption
    r.font.size = Pt(13)
    r.font.italic = True
    r.font.color.rgb = GRAY
    r.font.name = FONT
    return sh


def connector(slide, x1, y1, x2, y2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = GREEN
    c.line.width = Pt(1)
    return c


def oval_number(slide, left, top, size, number, fill=GREEN, color=WHITE):
    sh = shape(slide, MSO_SHAPE.OVAL, left, top, size, size, fill=fill, line=None)
    shape_text(sh, str(number), size=18, bold=True, color=color)
    return sh


def chevron(slide, left, top, width, height, text):
    sh = shape(slide, MSO_SHAPE.CHEVRON, left, top, width, height, fill=GREEN, line=None)
    shape_text(sh, text, size=13, bold=True, color=WHITE)
    return sh


def callout(slide, text, left, top, width, height):
    """Caja crema con borde verde grueso a la izquierda + comilla, texto en cursiva."""
    box = shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, fill=CREAM, line=None)
    shape(slide, MSO_SHAPE.RECTANGLE, left, top, 0.10, height, fill=GREEN, line=None)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(16)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)
    p = tf.paragraphs[0]
    rq = p.add_run()
    rq.text = '“ '
    rq.font.size = Pt(24)
    rq.font.bold = True
    rq.font.color.rgb = GREEN
    rq.font.name = FONT
    _add_rich_runs(p, text, 14, base_color=DARK, kw_color=GREEN, base_bold=False)
    return box


# ----------------------------- Contenido por slide -----------------------------
def build(prs):
    # ---- Slide 1: Portada ----
    s = _blank(prs)
    add_title(s, 'Defensa de Pasantía')
    tb = s.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(11.3), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'Propuesta de Sistema de Gestión Documental Digital para la Optimización de los Procesos Administrativos en la Gerencia Administrativa de la empresa Servicios y Suministros GEO, C.A.'
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = GREEN
    r.font.name = FONT
    tb2 = s.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(11.3), Inches(1.6))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate([
        'Defensa de Informe de Pasantías — Técnico Superior Universitario en Informática',
        'Pasante: Juan Moya (C.I. V-32.383.863)',
        'IUTECP "Elías Calixto Pompa" · El Tigre, Anzoátegui',
        'Servicios y Suministros GEO, C.A. · Agosto 2026',
    ]):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        r.font.size = Pt(15 if i == 0 else 13)
        r.font.color.rgb = DARK if i == 0 else GRAY
        r.font.bold = (i == 0)
        r.font.name = FONT
    # placeholder logo
    placeholder(s, Inches(5.17), Inches(5.9), Inches(3.0), Inches(0.9), 'Logo IUTECP')
    add_footer(s, 1)

    # ---- Slide 2: Introducción / ruta ----
    s = _blank(prs)
    add_title(s, 'Introducción', '00 · INTRODUCCIÓN')
    rutas = ['Realidad organizacional', 'Diagnóstico y objetivos', 'Marco teórico', 'Actividades (10 semanas)', 'Conclusiones']
    x = MARGIN
    bw = 2.15
    gap = (CONTENT_W - 5 * bw) / 4
    for i, t in enumerate(rutas):
        pill(s, x, 1.5, bw, 0.7, t)
        if i < 4:
            arr = s.shapes.add_textbox(Inches(x + bw + 0.02), Inches(1.5), Inches(gap), Inches(0.7))
            arr.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            rr = arr.text_frame.paragraphs[0].add_run()
            rr.text = '→'
            rr.font.size = Pt(20)
            rr.font.bold = True
            rr.font.color.rgb = GREEN
            rr.font.name = FONT
        x += bw + gap
    add_bullets(s, [
        'Buenos días, miembros del jurado, profesores y tutor académico.',
        'Presento y defiendo el informe de pasantías desarrollado entre **junio y agosto de 2026**.',
        'Recorrido integrado: la realidad de la empresa, el diagnóstico y objetivos, el marco que lo sustenta, las actividades ejecutadas y las conclusiones.',
    ], top=2.7)
    add_footer(s, 2)

    # ---- Slide 3: La empresa ----
    s = _blank(prs)
    add_title(s, 'La empresa: GEO, C.A.', '01 · REALIDAD ORGANIZACIONAL')
    pills = [('Sector: Petrolero', 1.35), ('Alcance: Falcón · Zulia', 1.35), ('Origen: 2016', 1.35)]
    x = MARGIN
    for t, _ in pills:
        pill(s, x, 1.35, 3.7, 0.6, t)
        x += 3.95
    add_bullets(s, [
        'Servicios y Suministros **GEO, C.A.** (RIF J-40833844-0), inscrita en el Registro Mercantil de Anzoátegui.',
        'Nació por iniciativa de Gustavo Ovalles para la industria, especialmente el **sector petrolero**.',
        'Oficinas en San José de Guanipa; ejecuta trabajos para empresas **petroleras y petroquímicas** del occidente.',
        'Portafolio: **mantenimiento industrial**, suministro de materiales y consultoría técnica.',
        'Misión, visión y valores (transparencia, honestidad, constancia) fueron la base ética de la propuesta.',
    ], top=2.2)
    add_footer(s, 3)

    # ---- Slide 4: Estructura organizacional ----
    s = _blank(prs)
    add_title(s, 'Estructura organizacional', '01 · REALIDAD ORGANIZACIONAL')
    # organigrama simple
    shape(s, MSO_SHAPE.RECTANGLE, 5.4, 1.4, 2.5, 0.6, fill=CREAM, line=GREEN)
    shape_text(s.shapes[-1], 'Presidencia', size=13)
    connector(s, 6.65, 2.0, 6.65, 2.3)
    shape(s, MSO_SHAPE.RECTANGLE, 5.0, 2.3, 3.3, 0.6, fill=CREAM, line=GREEN)
    shape_text(s.shapes[-1], 'Gerencia General', size=13)
    connector(s, 6.65, 2.9, 6.65, 3.2)
    # row of units
    units = ['Gerencia\nAdministrativa', 'Tecnología\nde la Información', 'Coordinaciones:\nCompras · Proyectos\nOperaciones · RRHH']
    x = 0.9
    uw = 3.5
    for u in units:
        shape(s, MSO_SHAPE.RECTANGLE, x, 3.2, uw, 1.0, fill=CREAM, line=GREEN)
        shape_text(s.shapes[-1], u, size=12)
        if x < 9:
            connector(s, x + uw, 3.7, x + uw + 0.45, 3.7)
        x += uw + 0.45
    add_bullets(s, [
        'Plantilla de **veinte trabajadores**.',
        'Mi pasantía se desarrolló en la **Gerencia de Administración**.',
    ], top=4.5)
    add_footer(s, 4)

    # ---- Slide 5: Gerencia de Administración ----
    s = _blank(prs)
    add_title(s, 'Gerencia de Administración', '01 · REALIDAD ORGANIZACIONAL')
    add_bullets(s, [
        'Conformada por el analista financiero, contabilidad, el asistente, **cuentas por pagar y tributación**, y cuentas por cobrar y facturación.',
        'Maneja los **procesos administrativos y financieros** de la organización.',
        'Incluye la elaboración y declaración de las **retenciones del IVA** ante el fisco nacional.',
    ], top=1.5)
    # caja resaltada
    box = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, 4.6, CONTENT_W, 1.2, fill=CREAM, line=GREEN, line_w=1.25)
    shape_text(box, 'Punto crítico: **GEO es contribuyente especial** → fallas en **retenciones de IVA** exponen a sanciones legales y fiscales.', size=15, bold=True, color=GREEN)
    add_footer(s, 5)

    # ---- Slide 6: Diagnóstico ----
    s = _blank(prs)
    add_title(s, 'Diagnóstico situacional', '02 · DIAGNÓSTICO Y OBJETIVOS')
    # gráfico mínimo: 3 archivos dispersos
    for i, lbl in enumerate(['Excel', 'PDF', 'Disco local']):
        sh = shape(s, MSO_SHAPE.RECTANGLE, 1.0 + i * 1.4, 1.5, 1.1, 0.8, fill=CREAM, line=GREEN)
        shape_text(sh, lbl, size=12)
    add_bullets(s, [
        'La **documentación contable y fiscal** se almacenaba de forma aislada en discos locales de los empleados.',
        'Dispersión → retrasos en el rastreo de comprobantes, fallas de versiones y riesgo de pérdida de datos confidenciales.',
        'Sin software estandarizado local: **duplicidad de funciones** y demoras en la **declaración de retenciones de IVA**.',
        'Un comprobante extraviado podía implicar **multa o improcedencia de deducción fiscal**.',
    ], top=2.7, height=2.6)
    callout(s, 'El problema no era de capacidad del personal, sino de **diseño del proceso**.', MARGIN, 5.45, CONTENT_W, 1.25)
    add_footer(s, 6)

    # ---- Slide 7: Ishikawa ----
    s = _blank(prs)
    add_title(s, 'Análisis con Ishikawa', '02 · DIAGNÓSTICO Y OBJETIVOS')
    if os.path.exists(RUTA_ISHIKAWA):
        s.shapes.add_picture(RUTA_ISHIKAWA, Inches(2.2), Inches(1.35), width=Inches(9.0))
    add_bullets(s, [
        'El diagrama de **Ishikawa** evidenció las causas raíz de la ineficiencia.',
        'La **dispersión de archivos** y la falta de un índice centralizado impactan directamente la eficiencia y el control fiscal.',
        'La solución apunta al **diseño del proceso documental**, no a la capacidad del personal.',
    ], top=5.6, height=1.4)
    add_footer(s, 7)

    # ---- Slide 8: Objetivos ----
    s = _blank(prs)
    add_title(s, 'Objetivos de la investigación', '02 · DIAGNÓSTICO Y OBJETIVOS')
    oval_number(s, MARGIN, 1.6, 0.7, 1)
    tb = s.shapes.add_textbox(Inches(1.4), Inches(1.55), Inches(10.5), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _add_rich_runs(p, 'General: proponer un **sistema de gestión documental digital** para optimizar los procesos administrativos de la gerencia.', 17, base_color=DARK, kw_color=GREEN, base_bold=True)
    espec = [
        'Diagnosticar la **situación procedimental** y los requerimientos de información.',
        'Evaluar el impacto de las deficiencias y **limitaciones de control interno** sobre la eficiencia y el **cumplimiento fiscal**.',
        'Diseñar la arquitectura lógica, los esquemas de datos y una demostración técnica funcional en **Rust**.',
    ]
    y = 2.7
    for i, e in enumerate(espec):
        oval_number(s, MARGIN, y, 0.55, i + 1)
        tb = s.shapes.add_textbox(Inches(1.25), Inches(y - 0.05), Inches(10.8), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _add_rich_runs(p, e, 15, base_color=DARK, kw_color=GREEN, base_bold=False)
        y += 1.15
    add_footer(s, 8)

    # ---- Slide 9: Metodología ----
    s = _blank(prs)
    add_title(s, 'Metodología', '02 · DIAGNÓSTICO Y OBJETIVOS')
    fases = ['Diagnóstico', 'Requerimientos', 'Diseño', 'Prototipo']
    cw = 2.7
    gap = (CONTENT_W - 4 * cw) / 3
    x = MARGIN
    for f in fases:
        chevron(s, x, 1.5, cw, 0.8, f)
        x += cw + gap
    add_bullets(s, [
        'Proyecto factible apoyado en **investigación de campo** y revisión documental.',
        'Fases de análisis de sistemas de **Kendall y Kendall**.',
        'Técnicas: observación participante, entrevistas no estructuradas, listas de cotejo, **matrices causa-efecto**.',
        'Diagrama de **Gantt** organizó las actividades a lo largo de las **diez semanas** reglamentarias.',
    ], top=2.7)
    placeholder(s, MARGIN, 5.4, CONTENT_W, 1.3, 'Espacio para imagen: Diagrama de Gantt (10 semanas)')
    add_footer(s, 9)

    # ---- Slide 10: Marco teórico ----
    s = _blank(prs)
    add_title(s, 'Marco teórico', '03 · MARCO TEÓRICO')
    card(s, MARGIN, 1.5, 3.9, 2.4, 'Gestión documental',
         'Administra documentos en soporte electrónico a lo largo de su ciclo de vida: **control de versiones**, **trazabilidad** y simplificación de procesos.')
    card(s, MARGIN + 4.2, 1.5, 3.9, 2.4, 'Tecnologías de información',
         'Arquitectura de software, bases de datos relacionales, **metadatos** (RIF, nº factura, tipo de retención) y **RBAC**.')
    card(s, MARGIN + 8.4, 1.5, 3.9, 2.4, 'Sistemas gerenciales',
         'Apoya la toma de decisiones; ingeniería de requerimientos y atributos de calidad, **mantenibilidad** y fiabilidad.')
    add_bullets(s, [
        'Capturar, clasificar, almacenar y disponer el documento sin rupturas en la **cadena de custodia**.',
    ], top=4.2)
    add_footer(s, 10)

    # ---- Slides 11-14: Etapas ----
    etapas = [
        ('Etapa I — Inducción y Diagnóstico (S1-2)', 'diagnostico',
         ['Inducción con el tutor industrial.', 'Inspección de facturas, órdenes de servicio y **comprobantes de retención de IVA**.', 'Análisis del procedimiento **SENIAT** (contribuyente especial).', 'Levantamiento del archivo digital: carpetas dispersas, **duplicidad**, riesgo de demora.'],
         'Espacio para imagen: Diagnóstico y Inspección de Archivos'),
        ('Etapa II — Levantamiento y Modelado (S3-4)', 'levantamiento',
         ['Mapeo de la secuencia documental en cuentas por cobrar, pagar y tributación.', 'Entrevistas con contador y analista financiero.', 'Definición de **metadatos clave** (RIF, nº factura, fecha, tipo de retención).', 'Modelado de diagramas de flujo (**Kendall y Kendall**) y ERS.'],
         'Espacio para imagen: Diagramas de Flujo ERS'),
        ('Etapa III — Diseño y Codificación en Rust (S5-7)', 'diseno',
         ['Políticas de **RBAC** conforme al Código Orgánico Tributario.', 'Rutas locales en **Linux Debian**.', 'Demostración modular en **Rust**: extracción automática de metadatos e indexación centralizada.', 'Filtrados inmediatos sin software comercial pesado.'],
         'Espacio para imagen: Arquitectura del Software y Demostración Técnica'),
        ('Etapa IV — Pruebas, Capacitación y Cierre (S8-10)', 'pruebas',
         ['Pruebas piloto con analista financiero y personal administrativo.', 'Evaluación del motor de búsqueda en **Rust** sobre volúmenes reales (**Linux Debian**).', 'Ajuste de algoritmos de filtrado con la retroalimentación.', 'Jornada de capacitación, manuales de usuario y documentación técnica.'],
         'Espacio para imagen: Pruebas de Usabilidad e Integración'),
    ]
    for idx, (title, _, bullets, cap) in enumerate(etapas, start=1):
        s = _blank(prs)
        add_title(s, title, '04 · ACTIVIDADES (10 SEMANAS)')
        # badge romano
        shape(s, MSO_SHAPE.OVAL, MARGIN, 1.35, 0.7, 0.7, fill=GREEN, line=None)
        shape_text(s.shapes[-1], ['I', 'II', 'III', 'IV'][idx - 1], size=18, bold=True, color=WHITE)
        # progreso: 4 puntos
        px = MARGIN + 1.0
        for i in range(4):
            fill = GREEN if i < idx else CREAM
            shape(s, MSO_SHAPE.OVAL, px + i * 0.45, 1.55, 0.28, 0.28, fill=fill, line=GREEN, line_w=0.75)
        add_bullets(s, [(b, 0) for b in bullets], top=2.4, height=3.0)
        placeholder(s, MARGIN, 5.5, CONTENT_W, 1.3, cap)
        add_footer(s, 10 + idx)

    # ---- Slide 15: Conclusiones ----
    s = _blank(prs)
    add_title(s, 'Conclusiones', '05 · CONCLUSIONES Y RECOMENDACIONES')
    add_bullets(s, [
        ('O1: el diagnóstico constató la inexistencia de **estructura digital centralizada** → duplicidad, demoras y vulnerabilidad tributaria.', 0),
        ('O2: el análisis causa-efecto demostró que la falta de flujo digital estandarizado compromete la productividad y eleva el riesgo de **sanciones**.', 0),
        ('O3: las pruebas piloto validaron la viabilidad técnica en **Rust**: tiempos de localización reducidos, interfaz intuitiva y **RBAC** seguro.', 0),
        'El plan de estudios aportó el marco organizacional, la rigurosidad metodológica y los fundamentos de modelado, programación y auditoría.',
    ], top=1.5, height=3.5)
    callout(s, 'Centralizar y **sistematizar la información** es la base para cumplir, auditar y crecer.', MARGIN, 5.3, CONTENT_W, 1.3)
    add_footer(s, 15)

    # ---- Slide 16: Recomendaciones ----
    s = _blank(prs)
    add_title(s, 'Recomendaciones', '05 · CONCLUSIONES Y RECOMENDACIONES')
    card(s, MARGIN, 1.5, 3.9, 4.0, 'A la empresa',
         '• Adopción progresiva del sistema.\n• Esquema de nombrado por **metadatos**.\n• Integrar en cuentas por cobrar y por pagar.\n• Capacitar al personal y abrir pasantías.')
    card(s, MARGIN + 4.2, 1.5, 3.9, 4.0, 'Al IUTECP',
         '• Talleres de **lenguajes y arquitecturas emergentes**.\n• Comunicación entre pasantías y tutores.\n• Proyectos integradores (auditoría y gestión documental).')
    card(s, MARGIN + 8.4, 1.5, 3.9, 4.0, 'A los futuros pasantes',
         '• Priorizar **diagnóstico de campo** y requerimientos desde las primeras semanas.\n• Mantener comunicación constante con ambos tutores.')
    add_footer(s, 16)

    # ---- Slide 17: Cierre ----
    s = _blank(prs)
    add_title(s, 'Cierre', '05 · CONCLUSIONES Y RECOMENDACIONES')
    tb = s.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.9), Inches(3.2))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = [
        'La pasantía integró los conocimientos de la carrera en un entorno operativo real.',
        'La demostración funcional en **Rust** centraliza el archivo, elimina brechas operacionales y garantiza la trazabilidad exigida en auditorías internas y tributarias.',
        'La sistematización de la información es pilar de la productividad y del fortalecimiento del control interno.',
    ]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        p.alignment = PP_ALIGN.CENTER
        _add_rich_runs(p, ln, 17, base_color=DARK, kw_color=GREEN, base_bold=False)
    callout(s, 'Incluso con recursos locales y un lenguaje de sistemas como **Rust**, es posible resolver una necesidad real de la industria nacional.', MARGIN, 5.2, CONTENT_W, 1.1)
    tb2 = s.shapes.add_textbox(Inches(1.2), Inches(6.5), Inches(10.9), Inches(0.5))
    p = tb2.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'Muchas gracias'
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = GREEN
    r.font.name = FONT
    add_footer(s, 17)


def main():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    build(prs)
    prs.save(RUTA_SALIDA)
    print('OK ->', RUTA_SALIDA, '| slides:', len(prs.slides._sldIdLst))


if __name__ == '__main__':
    main()
