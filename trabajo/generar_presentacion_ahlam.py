# -*- coding: utf-8 -*-
"""Genera presentacion_ahlam.pptx (13 diapositivas, defensa de pasantía de Ahlam Alrifaai Alrifaie).

Estilo: Editorial Académico Corporativo (minimalista, sobrio).
Fondo marfil menta #F4F7F5, texto pizarra #1E293B, acento azul marino #0F172A,
resalte botánico #2D5A47, bordes ultrafinos #E2E8F0, sombras casi invisibles.
"""
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# ----------------------------- Paleta -----------------------------
IVORY = RGBColor(0xF4, 0xF7, 0xF5)      # Marfil Menta (fondo)
SLATE = RGBColor(0x1E, 0x29, 0x3B)       # Pizarra Oscuro (texto)
NAVY = RGBColor(0x0F, 0x17, 0x2A)       # Azul Marino Profundo (acento)
BOTANICAL = RGBColor(0x2D, 0x5A, 0x47)  # Verde Botánico (resalte)
BOTANICAL_TINT = RGBColor(0xEA, 0xF2, 0xEE)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)     # Borde ultrafino
ALERT = RGBColor(0xFE, 0xF3, 0xC7)      # Pastel cálido (alerta sísmica)
GRAY = RGBColor(0x64, 0x74, 0x8B)
DIV = RGBColor(0xCB, 0xD5, 0xE1)        # Divisoria sutil
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

BODY_FONT = 'Inter'
TITLE_FONT = 'Inter'

RUTA_SALIDA = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'presentacion_ahlam.pptx')

W = 13.333
H = 7.5
MARGIN = 0.67
CONTENT_W = W - 2 * MARGIN
TOTAL_SLIDES = 13
AUTOR = 'Ahlam Alrifaai Alrifaie'


# ----------------------------- Helpers -----------------------------
def _blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = IVORY
    return slide


def _add_rich_runs(p, text, size, base_color=SLATE, kw_color=NAVY, base_bold=False):
    parts = re.split(r'\*\*(.+?)\*\*', text)
    for i, part in enumerate(parts):
        if not part:
            continue
        r = p.add_run()
        r.text = part
        r.font.size = Pt(size)
        r.font.name = BODY_FONT
        if i % 2 == 1:
            r.font.bold = True
            r.font.color.rgb = kw_color
        else:
            r.font.bold = base_bold
            r.font.color.rgb = base_color


def add_soft_shadow(shape, blur=8, offset=2, color=RGBColor(0, 0, 0), alpha=0.03):
    sp_pr = shape._element.spPr
    existing = sp_pr.find(qn('a:effectLst'))
    if existing is not None:
        sp_pr.remove(existing)
    effect = parse_xml(
        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:outerShdw blurRad="{blur}" dist="{dist}" dir="5400000" rotWithShape="0">'
        '<a:srgbClr val="{color}"><a:alpha val="{alpha}"/></a:srgbClr>'
        '</a:outerShdw></a:effectLst>'.format(
            blur=int(blur * 12700), dist=int(offset * 12700),
            color='%02X%02X%02X' % (color[0], color[1], color[2]), alpha=int(alpha * 100000)))
    sp_pr.append(effect)


def shape(slide, kind, left, top, width, height, fill=None, line=BORDER, line_w=1, dash=False):
    sh = slide.shapes.add_shape(kind, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
        if dash:
            sh.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    sh.shadow.inherit = False
    return sh


def rounded_card(slide, left, top, width, height, fill=WHITE, line=BORDER, line_w=1, radius_in=0.1, shadow=True):
    sh = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill=fill, line=line, line_w=line_w)
    try:
        sh.adjustments[0] = min(0.5, radius_in / height)
    except Exception:
        pass
    if shadow:
        add_soft_shadow(sh)
    return sh


def card_text(sh, title, body, title_size=15, body_size=12, title_color=NAVY, body_color=SLATE):
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(8)
    tf.margin_left = Pt(11)
    tf.margin_right = Pt(11)
    tf.margin_bottom = Pt(8)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(title_size)
    r.font.bold = True
    r.font.color.rgb = title_color
    r.font.name = TITLE_FONT
    if body:
        p2 = tf.add_paragraph()
        _add_rich_runs(p2, body, body_size, base_color=body_color, kw_color=NAVY, base_bold=False)


def card(slide, left, top, width, height, title, body, fill=WHITE, title_size=15, body_size=12,
         title_color=NAVY, body_color=SLATE, line_w=1, **kw):
    sh = rounded_card(slide, left, top, width, height, fill=fill, line_w=line_w, **kw)
    card_text(sh, title, body, title_size=title_size, body_size=body_size,
              title_color=title_color, body_color=body_color)
    return sh


def chip(slide, left, top, width, text, fill=BOTANICAL_TINT, text_color=BOTANICAL, size=11):
    sh = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, 0.4, fill=fill, line=None)
    try:
        sh.adjustments[0] = 0.38 / 0.4
    except Exception:
        pass
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _add_rich_runs(p, text, size, base_color=text_color, kw_color=BOTANICAL, base_bold=True)
    return sh


def placeholder(slide, left, top, width, height, caption):
    sh = shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, fill=RGBColor(0xF1, 0xF5, 0xF4),
               line=BORDER, line_w=1.25, dash=True)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = caption
    r.font.size = Pt(12)
    r.font.italic = True
    r.font.color.rgb = GRAY
    r.font.name = BODY_FONT
    return sh


def add_title(slide, text, kicker=None):
    if kicker:
        kt = slide.shapes.add_textbox(Inches(MARGIN), Inches(0.45), Inches(CONTENT_W - 1.0), Inches(0.3))
        kp = kt.text_frame.paragraphs[0]
        kr = kp.add_run()
        kr.text = kicker.upper()
        kr.font.size = Pt(12)
        kr.font.bold = True
        kr.font.color.rgb = NAVY
        kr.font.name = BODY_FONT
        top_title = 0.82
        top_line = 1.42
    else:
        top_title = 0.6
        top_line = 1.25
    tb = slide.shapes.add_textbox(Inches(MARGIN), Inches(top_title), Inches(CONTENT_W - 1.0), Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(25)
    r.font.bold = True
    r.font.color.rgb = NAVY
    r.font.name = TITLE_FONT
    line = shape(slide, MSO_SHAPE.RECTANGLE, MARGIN, top_line, CONTENT_W, 0.02, fill=BORDER, line=None)
    return tb


def add_footer(slide, n):
    shape(slide, MSO_SHAPE.RECTANGLE, MARGIN, 6.92, CONTENT_W, 0.02, fill=BORDER, line=None)
    tb = slide.shapes.add_textbox(Inches(MARGIN), Inches(7.0), Inches(9), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = 'IUTSO 2026 · ' + AUTOR + ' · Defensa de Pasantía'
    r.font.size = Pt(9)
    r.font.color.rgb = GRAY
    r.font.name = BODY_FONT
    nb = slide.shapes.add_textbox(Inches(W - 1.2), Inches(7.0), Inches(0.8), Inches(0.3))
    p = nb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f'{n:02d}'
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = NAVY
    r.font.name = BODY_FONT


def add_bullets(slide, items, top=1.7, left=MARGIN, width=CONTENT_W, height=4.0, marker='▪  '):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        text, level = it if isinstance(it, tuple) else (it, 0)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(10)
        p.level = level
        rb = p.add_run()
        rb.text = ('        ' + marker) if level > 0 else marker
        rb.font.color.rgb = BOTANICAL
        rb.font.bold = True
        rb.font.size = Pt(15 if level == 0 else 13)
        rb.font.name = BODY_FONT
        _add_rich_runs(p, text, 15 if level == 0 else 13, base_color=SLATE, kw_color=NAVY, base_bold=False)


def add_arrow(slide, x, y, size=26):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(1.0), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = '➔'
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = NAVY
    r.font.name = BODY_FONT
    return tb


# ----------------------------- Diapositivas -----------------------------
def slide_portada(prs):
    s = _blank(prs)
    placeholder(s, MARGIN, 0.4, 1.7, 1.2, '[LOGO IUTSO]')
    mb = s.shapes.add_textbox(Inches(2.6), Inches(0.35), Inches(10.0), Inches(1.3))
    mb.text_frame.word_wrap = True
    membrete = [
        'REPÚBLICA BOLIVARIANA DE VENEZUELA',
        'MINISTERIO DEL PODER POPULAR PARA LA EDUCACIÓN UNIVERSITARIA',
        'INSTITUTO UNIVERSITARIO DE TECNOLOGÍA SUPERIOR DE ORIENTE (IUTSO)',
        'EL TIGRE, ESTADO ANZOÁTEGUI',
    ]
    first = True
    for line in membrete:
        p = mb.text_frame.paragraphs[0] if first else mb.text_frame.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        r.font.size = Pt(12.5)
        r.font.bold = True
        r.font.color.rgb = SLATE
        r.font.name = BODY_FONT
    shape(s, MSO_SHAPE.RECTANGLE, MARGIN, 1.82, CONTENT_W, 0.025, fill=DIV, line=None)
    sh = rounded_card(s, 1.4, 2.5, 10.53, 2.4, fill=WHITE, line=BORDER, line_w=1, radius_in=0.1)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(16)
    tf.margin_right = Pt(16)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'DESIGN AND WRITING OF A BILINGUAL (SPANISH-ENGLISH) CORPORATE BROCHURE FOR THE MARKETING AND COMMUNICATIONS DEPARTMENT AT GRUPO SOL Y SOMBRA, C.A., LOCATED IN EL TIGRE, ANZOÁTEGUI STATE'
    r.font.size = Pt(19)
    r.font.bold = True
    r.font.color.rgb = NAVY
    r.font.name = TITLE_FONT
    shape(s, MSO_SHAPE.RECTANGLE, MARGIN, 6.35, CONTENT_W, 0.025, fill=DIV, line=None)
    fb = s.shapes.add_textbox(Inches(MARGIN), Inches(6.5), Inches(CONTENT_W), Inches(0.5))
    p = fb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'PRESENTATION BY: ' + AUTOR
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = NAVY
    r.font.name = BODY_FONT


def slide_empresa(prs):
    s = _blank(prs)
    add_title(s, 'Grupo Sol y Sombra, C.A.', '01. LA ORGANIZACIÓN')
    col_w = (CONTENT_W - 0.4) / 2
    card(s, MARGIN, 1.7, col_w, 4.6, 'Resumen Ejecutivo',
         'Empresa consolidada en **El Tigre, Estado Anzoátegui**, bajo el modelo de negocio '
         '**"Calidad Económica": variedad de inventario con precios accesibles para el mercado regional.',
         title_size=17, body_size=14)
    rx = MARGIN + col_w + 0.4
    areas = [
        'Electrodomésticos y Tecnología',
        'Artículos para el Hogar',
        'Herramientas de Ferretería',
    ]
    for i, a in enumerate(areas):
        card(s, rx, 1.7 + i * 1.6, col_w, 1.4, a, '', title_size=15, body_size=12, title_color=BOTANICAL)
    add_footer(s, 2)


def slide_ubicacion(prs):
    s = _blank(prs)
    add_title(s, 'Ubicación Geográfica', '02. CONTEXTO')
    col_w = (CONTENT_W - 0.4) / 2
    card(s, MARGIN, 1.7, col_w, 4.6, 'Dirección Física',
         '▪  Calle Bolívar\n▪  Edificio Sol y Sombra\n▪  Local 34\n▪  Zona Centro\n'
         '▪  El Tigre\n▪  Estado Anzoátegui', body_size=15, title_size=17)
    placeholder(s, MARGIN + col_w + 0.4, 1.7, col_w, 4.6, '[MAPA / UBICACIÓN GEOGRÁFICA]')
    add_footer(s, 3)


def slide_estructura(prs):
    s = _blank(prs)
    add_title(s, 'Organigrama Departamental', '03. ESTRUCTURA')
    cw = 3.0
    y = 2.7
    h = 1.6
    x1 = MARGIN
    x2 = MARGIN + cw + 1.0
    x3 = MARGIN + 2 * (cw + 1.0)
    card(s, x1, y, cw, h, 'Gerencia de Ventas', 'Dirección general de operaciones comerciales.',
         title_size=15, body_size=12)
    add_arrow(s, x1 + cw + 0.05, y + 0.45)
    card(s, x2, y, cw, h, 'Mercadeo y Comunicaciones', 'Imagen y promoción institucional.',
         title_size=15, body_size=12)
    add_arrow(s, x2 + cw + 0.05, y + 0.45)
    sh = rounded_card(s, x3, y, cw, h, fill=BOTANICAL_TINT, line=BOTANICAL, line_w=1.25, radius_in=0.1)
    card_text(sh, 'Pasante de Inglés', 'Asistente de Contenido Bilingüe.', title_size=15, body_size=12,
              title_color=BOTANICAL, body_color=SLATE)
    add_footer(s, 4)


def slide_problema(prs):
    s = _blank(prs)
    add_title(s, 'Identificación del Problema', '04. DIAGNÓSTICO')
    sh = rounded_card(s, MARGIN, 1.8, CONTENT_W, 4.0, fill=WHITE, line=BORDER, line_w=1, radius_in=0.1)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(18)
    tf.margin_right = Pt(18)
    p = tf.paragraphs[0]
    _add_rich_runs(p,
        'La empresa carecía de una **herramienta promocional bilingüe formal** (información solo en español y '
        'de forma verbal). Esto generaba **barreras de comunicación** en un mercado activo como El Tigre, '
        'limitando su proyección e interacción internacional.', 17, base_color=SLATE, kw_color=NAVY, base_bold=False)
    add_footer(s, 5)


def slide_objetivos(prs):
    s = _blank(prs)
    add_title(s, 'Objetivos del Proyecto', '05. PROPÓSITO')
    card(s, MARGIN, 1.7, CONTENT_W, 1.1, 'Objetivo General',
         'Crear un **folleto / brochure bilingüe** para el Departamento de Mercadeo de Grupo Sol y Sombra, C.A.',
         title_size=17, body_size=13)
    specs = [
        ('1. Analizar', 'Información institucional y catálogo existente.'),
        ('2. Diseñar', 'Diagramación y contenido bilingüe (ES/EN).'),
        ('3. Desarrollar y Entregar', 'Producto final digital listo para su uso.'),
    ]
    cw = (CONTENT_W - 0.6) / 3
    for i, (t, b) in enumerate(specs):
        card(s, MARGIN + i * (cw + 0.3), 3.1, cw, 3.2, t, b, title_size=16, body_size=13)
    add_footer(s, 6)


def slide_plan(prs):
    s = _blank(prs)
    add_title(s, 'Plan de Actividades', '06. CRONOGRAMA')
    semanas = [
        ('01', '06/01 - 06/05', 'Inducción y evaluación de material publicitario.'),
        ('02', '06/08 - 06/12', 'Investigación corporativa y análisis del mercado comercial.'),
        ('03', '06/15 - 06/19', 'Recolección de fichas técnicas de productos de hogar y tecnología.'),
        ('04-05', '06/22 - 07/03', 'Maquetación en español, traducción a inglés y ensamblaje digital. (Pausa de 3 días por emergencia sísmica nacional; plan de recuperación aplicado).'),
        ('06', '07/06 - 07/10', 'Presentación del borrador al tutor institucional para correcciones.'),
        ('07', '07/13 - 07/15', 'Cumplimiento del plan de recuperación de días y entrega oficial.'),
    ]
    seg_w = CONTENT_W / len(semanas)
    line_y = 2.7
    shape(s, MSO_SHAPE.RECTANGLE, MARGIN, line_y, CONTENT_W, 0.03, fill=BORDER, line=None)
    for i, (sem, fechas, desc) in enumerate(semanas):
        xc = MARGIN + seg_w * i + seg_w / 2
        node = ALERT if sem == '04-05' else NAVY
        sh = shape(s, MSO_SHAPE.OVAL, xc - 0.15, line_y - 0.13, 0.3, 0.3, fill=node, line=WHITE, line_w=1.5)
        fill = ALERT if sem == '04-05' else WHITE
        cw = seg_w - 0.2
        card(s, MARGIN + seg_w * i + 0.1, line_y + 0.35, cw, 3.4, 'Sem ' + sem, fechas + '\n' + desc,
             fill=fill, title_size=13, body_size=10.5, line_w=1)
    add_footer(s, 7)


def slide_vinculacion(prs):
    s = _blank(prs)
    add_title(s, 'Vinculación con el Perfil Profesional', '07. PERFIL')
    bloques = [
        ('Inglés Técnico', 'Precisión gramatical en contexto comercial.'),
        ('Traducción', 'Dominio de terminología técnica de productos.'),
        ('Redacción', 'Estructuración de contenidos corporativos.'),
    ]
    cw = (CONTENT_W - 0.6) / 3
    for i, (t, b) in enumerate(bloques):
        card(s, MARGIN + i * (cw + 0.3), 1.8, cw, 3.8, t, b, title_size=16, body_size=13, title_color=BOTANICAL)
        chip(s, MARGIN + i * (cw + 0.3) + 0.2, 4.9, cw - 0.4, t, fill=BOTANICAL_TINT, text_color=BOTANICAL, size=11)
    add_footer(s, 8)


def slide_conocimientos(prs):
    s = _blank(prs)
    add_title(s, 'Conocimientos Teóricos y Prácticos', '08. APRENDIZAJE')
    col_w = (CONTENT_W - 0.4) / 2
    card(s, MARGIN, 1.7, col_w, 4.6, 'Teóricos',
         '✓  Vocabulario técnico en **Inglés para Fines Específicos (ESP)**.\n'
         '✓  Normas de comunicación corporativa bilingüe.', body_size=15, title_size=17)
    card(s, MARGIN + col_w + 0.4, 1.7, col_w, 4.6, 'Prácticos',
         '✓  Maquetación digital de folletos.\n✓  Trabajo en equipo corporativo.\n'
         '✓  Resolución de problemas en proyectos reales.', body_size=15, title_size=17)
    add_footer(s, 9)


def slide_demostracion(prs):
    s = _blank(prs)
    add_title(s, 'Demostración del Producto Final', '09. EVIDENCIAS')
    caps = ['[PÁGINA BROCHURE ESPAÑOL]', '[PÁGINA BROCHURE INGLÉS]', '[EVIDENCIA FOTOGRÁFICA]']
    cw = (CONTENT_W - 0.8) / 3
    for i, c in enumerate(caps):
        left = MARGIN + i * (cw + 0.4)
        mk = rounded_card(s, left, 1.8, cw, 4.0, fill=RGBColor(0xF1, 0xF5, 0xF4), line=BORDER, line_w=1.25, radius_in=0.08)
        mk.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        tf = mk.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = c
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = SLATE
        r.font.name = BODY_FONT
    add_footer(s, 10)


def slide_conclusiones(prs):
    s = _blank(prs)
    add_title(s, 'Conclusiones', '10. CIERRE')
    sh = rounded_card(s, 1.6, 1.9, CONTENT_W - 3.2, 3.6, fill=WHITE, line=BORDER, line_w=1, radius_in=0.1)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(18)
    tf.margin_right = Pt(18)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = '✔  '
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = BOTANICAL; r.font.name = BODY_FONT
    r2 = p.add_run()
    r2.text = ('Cumplimiento exitoso del objetivo al dotar a la empresa de una **herramienta bilingüe** '
               'que elimina barreras idiomáticas y fortalece su imagen corporativa.')
    r2.font.size = Pt(17); r2.font.color.rgb = SLATE; r2.font.name = BODY_FONT
    add_footer(s, 11)


def slide_recomendaciones(prs):
    s = _blank(prs)
    add_title(s, 'Recomendaciones Estratégicas', '11. SUGERENCIAS')
    recs = [
        ('A la Empresa', 'Mantener actualizado el brochure digital en todas sus plataformas.'),
        ('Al IUTSO', 'Continuar ofreciendo talleres sobre herramientas digitales de traducción y diseño.'),
        ('A Futuros Pasantes', 'Administrar el tiempo rigurosamente desde el primer día.'),
    ]
    cw = (CONTENT_W - 0.6) / 3
    for i, (t, b) in enumerate(recs):
        card(s, MARGIN + i * (cw + 0.3), 1.8, cw, 3.8, t, b, title_size=16, body_size=13)
    add_footer(s, 12)


def slide_cierre(prs):
    s = _blank(prs)
    kt = s.shapes.add_textbox(Inches(MARGIN), Inches(0.45), Inches(CONTENT_W - 1.0), Inches(0.3))
    kp = kt.text_frame.paragraphs[0]
    kr = kp.add_run()
    kr.text = '12. FINAL'
    kr.font.size = Pt(12)
    kr.font.bold = True
    kr.font.color.rgb = NAVY
    kr.font.name = BODY_FONT
    sh = rounded_card(s, 1.6, 2.2, CONTENT_W - 3.2, 3.0, fill=WHITE, line=BORDER, line_w=1, radius_in=0.12)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = '¡Muchas Gracias por su Atención!'
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = TITLE_FONT
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run(); r.text = 'Thank you very much for your attention'
    r.font.size = Pt(17); r.font.italic = True; r.font.color.rgb = BOTANICAL; r.font.name = BODY_FONT
    fb = s.shapes.add_textbox(Inches(MARGIN), Inches(6.7), Inches(CONTENT_W), Inches(0.4))
    p = fb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = 'IUTSO 2026 · ' + AUTOR
    r.font.size = Pt(12); r.font.color.rgb = GRAY; r.font.name = BODY_FONT


# ----------------------------- Main -----------------------------
def main():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide_portada(prs)
    slide_empresa(prs)
    slide_ubicacion(prs)
    slide_estructura(prs)
    slide_problema(prs)
    slide_objetivos(prs)
    slide_plan(prs)
    slide_vinculacion(prs)
    slide_conocimientos(prs)
    slide_demostracion(prs)
    slide_conclusiones(prs)
    slide_recomendaciones(prs)
    slide_cierre(prs)
    prs.save(RUTA_SALIDA)
    print('Guardado:', RUTA_SALIDA, '| diapositivas:', len(prs.slides._sldIdLst))


if __name__ == '__main__':
    main()
